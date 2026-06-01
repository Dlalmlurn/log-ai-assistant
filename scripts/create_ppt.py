#!/usr/bin/env python3
"""Generate project stage PPT — UEBA anomaly detection system.
Design principles:
  - Diagrams over bullet lists
  - Technical rationale over naming
  - Clean visual hierarchy
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Colors ────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x11, 0x1A)   # near-black
CARD    = RGBColor(0x16, 0x1B, 0x26)   # surface
BORDER  = RGBColor(0x25, 0x2E, 0x3F)   # subtle border
CYAN    = RGBColor(0x38, 0xBD, 0xF8)   # primary accent
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)   # secondary
GREEN   = RGBColor(0x4A, 0xDE, 0x80)   # success
ORANGE  = RGBColor(0xFB, 0x92, 0x3C)   # warn
RED     = RGBColor(0xF8, 0x71, 0x71)   # danger
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
DIM     = RGBColor(0x64, 0x6A, 0x76)
SOFT    = RGBColor(0xD1, 0xD5, 0xDB)
SUBTLE  = RGBColor(0x1F, 0x28, 0x37)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Reusable builders ─────────────────────────────────────────

def _solid_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h, fill=CARD, border=None, radius=None):
    """Rounded rectangle card. Returns shape."""
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border; shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _txt(slide, x, y, w, h, text, size=13, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, name="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
    """Single-line text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    tb.text_frame.auto_size = None
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tb.text_frame


def _para(tf, text, size=12, color=WHITE, bold=False, space_after=Pt(4)):
    """Append paragraph to existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.space_after = space_after
    return p


def _line(slide, x, y, w, color=CYAN, h=Pt(2.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def _arrow(slide, x1, y1, x2, y2, color=DIM, width=Pt(1.5)):
    """Draw a connector line with arrowhead."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    tail = etree.SubElement(conn.line._ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'sm')
    tail.set('len', 'sm')
    return conn


def _circle(slide, x, y, d, fill=CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s


def _node(slide, x, y, w, h, label, sub="", fill=CARD, accent=CYAN):
    """A labeled diagram node with accent left border."""
    body = _box(slide, x, y, w, h, fill=fill, border=BORDER, radius=0.08)
    _line(slide, x, y, Pt(4), accent, h)
    _txt(slide, x + Inches(0.2), y + Inches(0.08), w - Inches(0.3), Inches(0.32),
         label, size=11.5, color=WHITE, bold=True)
    if sub:
        _txt(slide, x + Inches(0.2), y + h - Inches(0.32), w - Inches(0.3), Inches(0.28),
             sub, size=9.5, color=DIM)
    return body


def _slide_title(slide, num, cn, en=""):
    """Slide header with number, Chinese title, English subtitle."""
    _txt(slide, Inches(0.7), Inches(0.4), Inches(0.5), Inches(0.4),
         f"{num:02d}", size=11, color=DIM, bold=True)
    _line(slide, Inches(1.3), Inches(0.55), Inches(1.2), CYAN)
    _txt(slide, Inches(1.3), Inches(0.32), Inches(8), Inches(0.45),
         cn, size=22, color=WHITE, bold=True)
    if en:
        _txt(slide, Inches(1.3), Inches(0.74), Inches(8), Inches(0.3),
             en, size=10, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)

# Left accent bar
_box(s, 0, 0, Inches(0.12), H, fill=CYAN)

# Brand mark — 3 dots
for i, clr in enumerate([CYAN, PURPLE, GREEN]):
    _circle(s, Inches(1.3) + i * Inches(0.35), Inches(1.8), Inches(0.18), fill=clr)

_txt(s, Inches(1.3), Inches(2.3), Inches(10), Inches(0.8),
     "UEBA 异常行为检测系统", size=44, color=WHITE, bold=True)
_txt(s, Inches(1.3), Inches(3.2), Inches(10), Inches(0.5),
     "从用户行为基线建模到 AI 智能研判的完整安全分析平台", size=16, color=GRAY)
# Tagline row
_line(s, Inches(1.3), Inches(3.9), Inches(2.5), CYAN)
tags = ["Docker 一体化", "自适应基线", "双引擎检测", "AI 研判"]
for i, t in enumerate(tags):
    _txt(s, Inches(1.3) + i * Inches(2.8), Inches(4.15), Inches(2.5), Inches(0.3),
         t, size=11, color=CYAN)

_txt(s, Inches(1.3), Inches(6.5), Inches(5), Inches(0.3),
     "2026-05-27   |   阶段汇报", size=11, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM & APPROACH  (why this matters)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 1, "我们要解决什么问题", "PROBLEM STATEMENT")

# Left: pain points
_txt(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4),
     "传统安全分析的瓶颈", size=15, color=ORANGE, bold=True)
pains = [
    ("日志割裂", "网络、终端、应用日志分散在各系统，\n安全分析师需要在多个平台间切换，无法关联分析。"),
    ("基线缺失", "没有用户\"正常行为\"的量化标准，\n检测只能依赖固定阈值或已知签名，\n未知威胁完全不可见。"),
    ("告警疲劳", "大量低质量告警缺乏优先级和上下文，\n分析师每天面对数百条告警却无从下手。"),
    ("研判耗时", "从告警到结论需要手动查询多个数据源，\n一次简单研判可能需要 30 分钟。"),
]
y = Inches(2.1)
for title, desc in pains:
    _box(s, Inches(0.7), y, Inches(5.5), Inches(1.05), fill=SUBTLE, border=BORDER, radius=0.06)
    _txt(s, Inches(1.1), y + Inches(0.1), Inches(2), Inches(0.28),
         title, size=13, color=ORANGE, bold=True)
    _txt(s, Inches(1.1), y + Inches(0.42), Inches(5), Inches(0.55),
         desc, size=10.5, color=GRAY)
    y += Inches(1.15)

# Right: our approach
_txt(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
     "我们的解决思路", size=15, color=CYAN, bold=True)

approaches = [
    ("归一化管道", "将所有日志统一解析为\n5W1H 结构，集中存储，\n打破数据孤岛。"),
    ("统计基线", "从历史行为中学习每个\n用户的正常模式，用量化\n偏差代替硬编码规则。"),
    ("多维加权", "6 个维度独立评估，\n自适应权重反映每个用户\n的行为稳定性差异。"),
    ("AI 研判", "检测到异常后自动调用\n大模型进行上下文分析，\n给出结论、证据和建议。"),
]
y = Inches(2.1)
for i, (title, desc) in enumerate(approaches):
    _box(s, Inches(7.0), y, Inches(5.5), Inches(1.05), fill=SUBTLE, border=BORDER, radius=0.06)
    _circle(s, Inches(7.3), y + Inches(0.15), Inches(0.22), fill=CYAN)
    _txt(s, Inches(7.3), y + Inches(0.17), Inches(0.22), Inches(0.22),
         str(i+1), size=10, color=BG, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, Inches(7.7), y + Inches(0.1), Inches(3), Inches(0.28),
         title, size=13, color=CYAN, bold=True)
    _txt(s, Inches(7.7), y + Inches(0.42), Inches(4.5), Inches(0.55),
         desc, size=10.5, color=GRAY)
    y += Inches(1.15)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — DATA PIPELINE FLOW DIAGRAM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 2, "数据如何流动", "PIPELINE ARCHITECTURE")

# Pipeline flow — horizontal
nodes = [
    ("日志生成", "模拟 VPN\n网关日志", GREEN),
    ("Filebeat", "采集投递", CYAN),
    ("Kafka", "消息缓冲\nraw_logs", ORANGE),
    ("Flink 流解析", "归一化 →\nparsed_logs", PURPLE),
    ("ClickHouse", "双路入库\n物化视图", CYAN),
    ("异常检测", "双引擎\n实时评分", RED),
    ("AI 研判", "大模型\n上下文分析", PURPLE),
]
nw = Inches(1.55)
nh = Inches(0.95)
gap = Inches(0.22)
start_x = Inches(0.7)
y_flow = Inches(1.7)

for i, (title, sub, accent) in enumerate(nodes):
    x = start_x + i * (nw + gap)
    _node(s, x, y_flow, nw, nh, title, sub, accent=accent)
    if i < len(nodes) - 1:
        ax = x + nw
        ay = y_flow + nh // 2
        _arrow(s, ax, ay, ax + gap, ay, color=DIM)

# Below: storage + consumer layer
_txt(s, Inches(0.7), Inches(3.1), Inches(3), Inches(0.3),
     "存储层", size=11, color=DIM, bold=True)
_line(s, Inches(0.7), Inches(3.35), Inches(12), BORDER, Pt(0.8))

storage_boxes = [
    ("security_logs", "主日志表\n38+字段·按月分区\n90天TTL"),
    ("ueba_user_baseline", "用户基线表\nEAV 窄表·T+1 更新\nmean/std/p50/p95/p99"),
    ("anomaly_events", "异常事件表\n含风险评分/偏离证据\nAI 研判状态"),
    ("ai_judgements", "AI 研判表\nattack_type/confidence\n关键证据/处置建议"),
]
bx_w = Inches(2.85)
bx_gap = Inches(0.22)
bx_start = Inches(0.7)
bx_y = Inches(3.55)
for i, (title, desc) in enumerate(storage_boxes):
    x = bx_start + i * (bx_w + bx_gap)
    _box(s, x, bx_y, bx_w, Inches(1.1), fill=SUBTLE, border=BORDER, radius=0.06)
    _txt(s, x + Inches(0.18), bx_y + Inches(0.1), bx_w - Inches(0.3), Inches(0.28),
         title, size=12, color=CYAN, bold=True)
    _txt(s, x + Inches(0.18), bx_y + Inches(0.42), bx_w - Inches(0.3), Inches(0.6),
         desc, size=10, color=GRAY)

# Consumer layer
_txt(s, Inches(0.7), Inches(5.05), Inches(3), Inches(0.3),
     "消费层", size=11, color=DIM, bold=True)
_line(s, Inches(0.7), Inches(5.3), Inches(12), BORDER, Pt(0.8))

consumers = [
    ("FastAPI Backend", "RESTful API\n11 个端点·结构化错误"),
    ("anomaly-detector", "后台常驻·实时消费\n批量写入·后台重试"),
    ("flink-submit daemon", "守护进程·自动监控\nFlink 作业存活"),
    ("React Frontend", "分析师工作台\n4 页面·15s 刷新"),
]
for i, (title, desc) in enumerate(consumers):
    x = bx_start + i * (bx_w + bx_gap)
    _box(s, x, Inches(5.5), bx_w, Inches(1.0), fill=SUBTLE, border=BORDER, radius=0.06)
    _txt(s, x + Inches(0.18), Inches(5.6), bx_w - Inches(0.3), Inches(0.25),
         title, size=11, color=GREEN, bold=True)
    _txt(s, x + Inches(0.18), Inches(5.88), bx_w - Inches(0.3), Inches(0.5),
         desc, size=9.5, color=GRAY)

# Key metric annotation
_txt(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.3),
     "端到端延迟 < 30s   |   5 个 Kafka Topic (3 分区)   |   默认 ~100 条/3.5s → ~1.1GB/天",
     size=10, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — DOCKER INTEGRATION (key emphasis — topology diagram)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 3, "Docker 一体化：一条命令构建整个平台", "ONE-COMMAND FULL-STACK DEPLOYMENT")

# One-command highlight
_highlight = _box(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.75), fill=RGBColor(0x1A, 0x25, 0x40))
_txt(s, Inches(1.2), Inches(1.45), Inches(10.5), Inches(0.55),
     "$ docker compose up -d     →     12 个服务自动编排 · 健康检查 · 依赖管理 · 零手动配置",
     size=19, color=CYAN, bold=True)

# Topology — 3 layers
layer_config = [
    ("基础设施层 (3)", [
        ("Kafka", "消息中枢 · KRaft 模式 · 5 个 Topic", CYAN),
        ("ClickHouse", "列存分析库 · 10 张表 · 物化视图", CYAN),
        ("Flink Cluster", "JobManager + TaskManager · 4 槽位", CYAN),
    ]),
    ("计算&采集层 (5)", [
        ("Filebeat", "日志采集 · JSON 编码 · gzip", PURPLE),
        ("log-generator", "模拟生成 · 双 profile · 确定性种子", PURPLE),
        ("flink-submit", "守护进程 · 作业监控 · 自动重提", PURPLE),
        ("anomaly-detector", "双引擎检测 · AI 研判触发", PURPLE),
        ("tester [test]", "pytest 套件 · 一键运行", PURPLE),
    ]),
    ("应用层 (4)", [
        ("backend (FastAPI)", "API 服务 · 11 端点 · 健康检查", GREEN),
        ("frontend (React)", "分析工作台 · 4 页面 · Vite 代理", GREEN),
        ("kafka-init", "Topic 初始化 · 一次性任务", GREEN),
        ("log-generator-scale", "高负载 profile · 50条/s", GREEN),
    ]),
]

y = Inches(2.35)
layer_accents = [CYAN, PURPLE, GREEN]
for li, (layer_name, services) in enumerate(layer_config):
    _txt(s, Inches(0.7), y, Inches(3), Inches(0.28),
         layer_name, size=11, color=layer_accents[li], bold=True)
    y += Inches(0.3)
    svc_w = Inches(3.7)
    svc_h = Inches(0.62)
    svc_gap = Inches(0.18)
    for j, (name, desc, clr) in enumerate(services):
        x = Inches(0.7) + j * (svc_w + svc_gap)
        _box(s, x, y, svc_w, svc_h, fill=SUBTLE, border=BORDER, radius=0.06)
        _txt(s, x + Inches(0.14), y + Inches(0.06), svc_w - Inches(0.2), Inches(0.26),
             name, size=11, color=clr, bold=True)
        _txt(s, x + Inches(0.14), y + Inches(0.34), svc_w - Inches(0.2), Inches(0.24),
             desc, size=8.5, color=DIM)
    y += Inches(0.82)

# Bottom: technical rationale
_txt(s, Inches(0.7), Inches(5.5), Inches(5.5), Inches(0.3),
     "为什么强调 Docker 一体化", size=13, color=CYAN, bold=True)
rationale = [
    "环境一致性: 开发/测试/演示完全相同的运行时，消除\"在我机器上能跑\"问题。",
    "依赖编排: healthcheck + depends_on 精确控制启动顺序，Kafka 就绪后才初始化 Topic，ClickHouse 健康后才启动 API。",
    "即时生效: 源码通过 read-only volume 挂载，修改 Python/TS 代码无需重建镜像，开发迭代秒级生效。",
    "容错自愈: flink-submit 守护进程监控 Flink 作业，JobManager 重启后自动重提，人工零干预。",
    "可扩展: 通过 profile 机制分离测试/压测/legacy 服务，核心服务不受影响。",
]
y = Inches(5.9)
for r in rationale:
    _txt(s, Inches(0.8), y, Inches(11.5), Inches(0.24),
         f"▸  {r}", size=10.5, color=GRAY)
    y += Inches(0.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — BASELINE: WHY AND HOW (key emphasis)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 4, "用户基线：从\"正常行为\"中学习异常", "BEHAVIOR BASELINE — CORE INNOVATION")

# Left: concept diagram — "before vs after"
_txt(s, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.35),
     "传统方式 vs 基线方式", size=14, color=CYAN, bold=True)

# Traditional
_old = _box(s, Inches(0.7), Inches(1.9), Inches(5.5), Inches(1.05), fill=RGBColor(0x2A, 0x1A, 0x1A), border=RED, radius=0.06)
_txt(s, Inches(1.1), Inches(2.0), Inches(2), Inches(0.25),
     "传统规则", size=13, color=RED, bold=True)
_txt(s, Inches(1.1), Inches(2.3), Inches(5), Inches(0.55),
     "if 失败次数 > 5: alert()\n固定阈值，对所有人都一样。\n张三日常失败率 30% 很正常，李四失败一次就异常。",
     size=10.5, color=GRAY)

# Baseline
_new = _box(s, Inches(0.7), Inches(3.1), Inches(5.5), Inches(1.65), fill=RGBColor(0x0A, 0x1E, 0x18), border=GREEN, radius=0.06)
_txt(s, Inches(1.1), Inches(3.2), Inches(2), Inches(0.25),
     "统计基线", size=13, color=GREEN, bold=True)
_txt(s, Inches(1.1), Inches(3.5), Inches(5), Inches(1.15),
     "if deviation(张三.登录失败率, 张三.历史正态分布) > 2σ: alert()\n"
     "张三的基线: μ=28%, σ=5% → 30% 在正常范围\n"
     "李四的基线: μ=2%, σ=1% → 5% 已经显著偏离\n\n"
     "同样的原始数据，不同的判断结论。这就是基线的本质。",
     size=10.5, color=GRAY)

# Right: how baseline is built
_txt(s, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.35),
     "基线是如何构建的 (T+1 模式)", size=14, color=CYAN, bold=True)

# Step diagram — vertical
steps = [
    ("每日聚合", "每天凌晨，从 security_logs 按 user_id\n聚合 23 个行为指标 → ueba_user_daily_features", CYAN),
    ("滑动窗口", "取最近 7 天的每日特征，对每个数值指标\n计算 mean / std / p50 / p95 / p99", PURPLE),
    ("置信度评估", "日覆盖度 × 样本量 × 特征覆盖率\n→ 置信度 < 0.3 自动降级到同类群组", ORANGE),
    ("写入基线", "以 EAV 窄表写入 ueba_user_baseline\n按 profile 分组为 5W1H 六类画像", GREEN),
]
y = Inches(1.9)
for i, (title, desc, clr) in enumerate(steps):
    _box(s, Inches(7.0), y, Inches(5.5), Inches(1.1), fill=SUBTLE, border=BORDER, radius=0.06)
    _circle(s, Inches(7.25), y + Inches(0.15), Inches(0.22), fill=clr)
    _txt(s, Inches(7.25), y + Inches(0.17), Inches(0.22), Inches(0.22),
         str(i+1), size=10, color=BG, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, Inches(7.65), y + Inches(0.1), Inches(2.5), Inches(0.25),
         title, size=12.5, color=clr, bold=True)
    _txt(s, Inches(7.65), y + Inches(0.4), Inches(4.6), Inches(0.6),
         desc, size=10, color=GRAY)
    # connector
    if i < len(steps) - 1:
        nx = Inches(7.0) + Inches(0.11)
        _arrow(s, nx, y + Inches(1.1), nx, y + Inches(1.15), color=DIM)
    y += Inches(1.25)

# Bottom insight
_insight = _box(s, Inches(7.0), Inches(6.45), Inches(5.5), Inches(0.65), fill=RGBColor(0x18, 0x1E, 0x2E))
_txt(s, Inches(7.3), Inches(6.52), Inches(5), Inches(0.5),
     "核心设计决策 (ADR-008): 基线必须来自历史统计，不能是生成器预设。\n"
     "这保证了检测的真实性——虚假基线等于没有基线。",
     size=10, color=ORANGE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — SIX-DIMENSION SCORING
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 5, "六维评分：每个用户都有不同的\"异常\"定义", "ADAPTIVE MULTI-DIMENSIONAL SCORING")

# Why 6 dimensions?
_txt(s, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.35),
     "为什么需要多维度", size=14, color=CYAN, bold=True)
_txt(s, Inches(0.7), Inches(1.75), Inches(5.5), Inches(0.8),
     "单一指标无法刻画\"异常\"。一个人凌晨3点登录，如果他是海外运维就完全正常，如果是朝九晚五的会计就高度可疑。"
     "\n需要同时从时间、地点、方式、频次等角度交叉验证，才能降低误报。",
     size=11, color=GRAY)

# Dimension radar — 6 boxes arranged 3x2
dims = [
    ("time", "时间维度", 0.15, "登录时刻是否在用户的活跃时段内", "基于小时直方图，计算左尾百分位。\n熵越高(24h均匀)→权重越低", CYAN),
    ("ip", "来源维度", 0.25, "来源 IP 是否在用户的已知集合中", "基于已知 IP 计数，1/(n+1) 概率。\n同/24子网减半惩罚", PURPLE),
    ("geo", "地域维度", 0.15, "地理位置是否与用户历史一致", "基于已知城市集合。\n新城市最低分 0.3，避免过度敏感", GREEN),
    ("access", "行为维度", 0.10, "操作/资源/UA 是否偏离日常", "匹配常用操作、资源、User-Agent。\n敏感资源 1.5x 加权", ORANGE),
    ("volume", "流量维度", 0.20, "单事件量级是否异常", "依赖窗口聚合 (已定义，延迟实现)。\n当前返回 None，不影响总分", RED),
    ("result", "结果维度", 0.15, "成功/失败率是否偏离历史", "对比历史登录成功率。\n极低或极高 → 高稳定性 → 高权重", PURPLE),
]

dim_w = Inches(3.75)
dim_h = Inches(1.35)
y = Inches(2.3)
for row in range(3):
    for col in range(2):
        idx = row * 2 + col
        dim_id, dim_cn, weight, question, method, clr = dims[idx]
        x = Inches(0.7) + col * (dim_w + Inches(0.25))
        yy = y + row * (dim_h + Inches(0.15))
        _box(s, x, yy, dim_w, dim_h, fill=SUBTLE, border=BORDER, radius=0.06)
        # accent top border
        _line(s, x, yy, dim_w, clr)
        _txt(s, x + Inches(0.18), yy + Inches(0.1), Inches(1), Inches(0.22),
             dim_id, size=10, color=clr, bold=True)
        _txt(s, x + Inches(0.9), yy + Inches(0.1), Inches(1.5), Inches(0.22),
             f"{dim_cn}  w={weight}", size=10, color=GRAY)
        _txt(s, x + Inches(0.18), yy + Inches(0.4), Inches(3.3), Inches(0.3),
             question, size=9.5, color=WHITE)
        _txt(s, x + Inches(0.18), yy + Inches(0.72), Inches(3.3), Inches(0.55),
             method, size=9, color=DIM)

# Bottom formula
_fbox = _box(s, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.55), fill=RGBColor(0x0F, 0x17, 0x22))
_txt(s, Inches(1.0), Inches(6.72), Inches(11.5), Inches(0.4),
     "自适应权重:  weight(dim) = normalize( 0.7 × stability_score(dim) + 0.3 )   →   组合分 = 0.5 × max_dev + 0.5 × Σ(w_i × dim_i)",
     size=11, color=CYAN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — DUAL ENGINE DETECTION FLOW
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 6, "双引擎协作：规则兜底 + 统计发现", "DUAL-ENGINE ARCHITECTURE")

# Flow diagram
_txt(s, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.3),
     "每条日志的处理路径", size=14, color=CYAN, bold=True)

# Incoming log
_log_node = _box(s, Inches(3.5), Inches(1.95), Inches(2.5), Inches(0.65), fill=CYAN, radius=0.06)
_log_node.fill.solid()
_txt(s, Inches(3.7), Inches(2.05), Inches(2.1), Inches(0.42),
     "NormalizedLog\n解析后的日志事件", size=10.5, color=BG, bold=True, align=PP_ALIGN.CENTER)

# Split to two paths
_arrow(s, Inches(4.75), Inches(2.6), Inches(2.3), Inches(3.15), color=DIM)
_arrow(s, Inches(4.75), Inches(2.6), Inches(8.0), Inches(3.15), color=DIM)

# Left: Rule Engine
_rule_box = _box(s, Inches(0.7), Inches(3.2), Inches(3.0), Inches(1.8), fill=SUBTLE, border=ORANGE, radius=0.06)
_txt(s, Inches(0.9), Inches(3.3), Inches(2.5), Inches(0.25),
     "规则引擎", size=13, color=ORANGE, bold=True)
_rule_items = [
    "IP 爆破 / 撞库 (滑动窗口)",
    "API 高频 / 敏感资源高频",
    "新 IP + 敏感操作 (组合)",
]
for i, ri in enumerate(_rule_items):
    _txt(s, Inches(0.9), Inches(3.65 + i*0.28), Inches(2.5), Inches(0.22),
         f"•  {ri}", size=9.5, color=GRAY)
_txt(s, Inches(0.9), Inches(4.6), Inches(2.5), Inches(0.28),
     "命中 → 直接高危 (score=85)", size=10, color=ORANGE, bold=True)

# Right: UEBA Scorer
_ueba_box = _box(s, Inches(6.0), Inches(3.2), Inches(6.5), Inches(1.8), fill=SUBTLE, border=PURPLE, radius=0.06)
_txt(s, Inches(6.2), Inches(3.3), Inches(3), Inches(0.25),
     "UEBA 评分器", size=13, color=PURPLE, bold=True)
_ueba_flow = [
    "1. 加载用户基线 (内存缓存, 5min 刷新)",
    "2. 六维度独立计算 surprise score (0-1)",
    "3. 自适应权重 × 维度分 → 加权组合",
    "4. ≥ 0.75 高危, ≥ 0.45 中危, 其余低危",
]
for i, uf in enumerate(_ueba_flow):
    _txt(s, Inches(6.2), Inches(3.65 + i*0.28), Inches(6), Inches(0.22),
         uf, size=9.5, color=GRAY)

# Merge logic
_merge = _box(s, Inches(3.0), Inches(5.4), Inches(4.5), Inches(0.9), fill=RGBColor(0x15, 0x20, 0x30), border=BORDER, radius=0.06)
_txt(s, Inches(3.2), Inches(5.48), Inches(4.1), Inches(0.25),
     "双引擎合并策略", size=12, color=CYAN, bold=True)
_txt(s, Inches(3.2), Inches(5.78), Inches(4.1), Inches(0.42),
     "规则命中 + UEBA 命中 → 规则告警吸收 UEBA 偏离作为佐证\n"
     "规则命中 only → 直接输出\n"
     "UEBA 命中 only → 输出统计偏离告警",
     size=10, color=GRAY)

_arrow(s, Inches(2.3), Inches(5.0), Inches(3.5), Inches(5.4), color=DIM)
_arrow(s, Inches(8.5), Inches(5.0), Inches(7.0), Inches(5.4), color=DIM)

# Output
_arrow(s, Inches(5.25), Inches(6.3), Inches(5.25), Inches(6.55), color=GREEN)
_out = _box(s, Inches(2.0), Inches(6.6), Inches(6.5), Inches(0.55), fill=GREEN, radius=0.06)
_out.fill.solid()
_txt(s, Inches(2.2), Inches(6.68), Inches(6.0), Inches(0.38),
     "AnomalyEvent  →  ClickHouse  →  AI 研判 (high-risk)",
     size=11, color=BG, bold=True, align=PP_ALIGN.CENTER)

# Right sidebar: key difference
_txt(s, Inches(7.5), Inches(1.4), Inches(5.0), Inches(0.28),
     "两者的分工", size=12, color=CYAN, bold=True)
_comp = [
    ("规则引擎", "已知攻击模式", "确定性匹配", "0 误报 (有把握)"),
    ("UEBA 评分", "未知行为偏离", "统计推断", "可发现新型威胁"),
]
for i, (engine, scope, method, advantage) in enumerate(_comp):
    _txt(s, Inches(7.5), Inches(1.78 + i*0.55), Inches(1.2), Inches(0.2),
         engine, size=10.5, color=CYAN if i == 0 else PURPLE, bold=True)
    _txt(s, Inches(8.8), Inches(1.78 + i*0.55), Inches(1.5), Inches(0.2),
         scope, size=10, color=WHITE)
    _txt(s, Inches(10.3), Inches(1.78 + i*0.55), Inches(1.5), Inches(0.2),
         method, size=10, color=GRAY)
    _txt(s, Inches(7.5), Inches(2.0 + i*0.55), Inches(5.0), Inches(0.2),
         advantage, size=9.5, color=GREEN if i == 1 else DIM)

_txt(s, Inches(7.5), Inches(3.0), Inches(5.0), Inches(0.6),
     "设计思路: 规则引擎覆盖\"已知已知\"，保证底线不丢；UEBA 评分探索\"已知未知\"，发现规则无法描述的异常。两者互补而非替代。",
     size=10, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — AI ANALYSIS (retry flow)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 7, "AI 研判：多模型容错与自动重试", "AI-POWERED TRIAGE WITH GRACEFUL DEGRADATION")

_txt(s, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.3),
     "设计目标：永远不出 mock 结果给用户", size=14, color=CYAN, bold=True)

# Retry flow diagram
# Entry
_trigger = _box(s, Inches(1.5), Inches(2.0), Inches(2.5), Inches(0.55), fill=CYAN, radius=0.06)
_trigger.fill.solid()
_txt(s, Inches(1.65), Inches(2.08), Inches(2.2), Inches(0.38),
     "高危异常触发 AI 研判", size=11, color=BG, bold=True, align=PP_ALIGN.CENTER)

# DeepSeek attempt
_arrow(s, Inches(2.75), Inches(2.55), Inches(3.2), Inches(3.1), color=DIM)
_ds = _box(s, Inches(3.2), Inches(3.15), Inches(4.0), Inches(1.0), fill=SUBTLE, border=CYAN, radius=0.06)
_txt(s, Inches(3.4), Inches(3.22), Inches(3.6), Inches(0.22),
     "DeepSeek API (主)", size=12, color=CYAN, bold=True)
_txt(s, Inches(3.4), Inches(3.5), Inches(3.6), Inches(0.55),
     "5 次重试, 指数退避\n1.5s → 3s → 6s → 12s\n成功 → 返回结果\n全部失败 → 降级到 DashScope",
     size=9.5, color=GRAY)

# DashScope fallback
_arrow(s, Inches(5.2), Inches(3.65), Inches(7.8), Inches(3.65), color=DIM)
_ds2 = _box(s, Inches(7.8), Inches(3.15), Inches(4.0), Inches(1.0), fill=SUBTLE, border=ORANGE, radius=0.06)
_txt(s, Inches(8.0), Inches(3.22), Inches(3.6), Inches(0.22),
     "DashScope (备用)", size=12, color=ORANGE, bold=True)
_txt(s, Inches(8.0), Inches(3.5), Inches(3.6), Inches(0.55),
     "同样 5 次重试\n成功 → 返回结果\n全部失败 → 保持 pending",
     size=9.5, color=GRAY)

# Pending state
_arrow(s, Inches(9.7), Inches(4.15), Inches(10.2), Inches(4.55), color=DIM)
_pending = _box(s, Inches(8.8), Inches(4.55), Inches(3.5), Inches(0.55), fill=RED, radius=0.06)
_pending.fill.solid()
_txt(s, Inches(8.95), Inches(4.63), Inches(3.2), Inches(0.38),
     "保持 pending 状态\n等待下一轮后台重试", size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Background retry
_retry = _box(s, Inches(8.8), Inches(5.4), Inches(3.5), Inches(0.9), fill=SUBTLE, border=GREEN, radius=0.06)
_txt(s, Inches(9.0), Inches(5.48), Inches(3.0), Inches(0.22),
     "后台恢复机制", size=12, color=GREEN, bold=True)
_txt(s, Inches(9.0), Inches(5.78), Inches(3.0), Inches(0.45),
     "启动时自动查询 db 中所有\npending 的 high-risk 异常\n逐条重新调用 AI 分析",
     size=9.5, color=GRAY)

_arrow(s, Inches(10.5), Inches(5.1), Inches(10.5), Inches(5.4), color=DIM)

# Right column: prompt engineering
_txt(s, Inches(0.7), Inches(4.55), Inches(7.5), Inches(0.3),
     "Prompt 工程：让模型输出可机器解析的结论", size=13, color=CYAN, bold=True)

prompt_parts = [
    ("系统角色", "\"你是企业安全分析助手。必须输出严格 JSON，不能输出额外文本。\""),
    ("输入上下文", "异常事件 (38+字段) + 用户5W1H基线画像 + 最多20条相关日志 + 窗口统计"),
    ("输出约束", "attack_type / risk_level / judgement / key_reasons / recommended_actions / confidence (0-1)"),
    ("容错处理", "LLM 可能返回 markdown 包裹的 JSON — 正则提取 + json.loads 二次尝试"),
]
y = Inches(4.95)
for title, desc in prompt_parts:
    _txt(s, Inches(0.9), y, Inches(1.5), Inches(0.22),
         title, size=10.5, color=CYAN, bold=True)
    _txt(s, Inches(2.5), y, Inches(5.5), Inches(0.22),
         desc, size=10, color=GRAY)
    y += Inches(0.28)

# Key insight
_txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.45),
     "核心思路: AI 调用是不可靠的——网络抖动、限流、模型繁忙都会失败。关键是 (1) 不丢数据 (pending 而非 failed) (2) 自动恢复 (启动时扫描未处理) (3) 优雅降级 (DeepSeek→DashScope 平滑切换)。",
     size=10, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — FRONTEND WORKBENCH
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 8, "前端设计：从数据到洞察的分析路径", "ANALYST WORKBENCH")

# User journey
_txt(s, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.3),
     "分析师的工作流程", size=14, color=CYAN, bold=True)

journey = [
    ("实时日志", "监控全局\n发现可疑", CYAN, "按用户/IP/时间筛选\n10s 自动刷新\n快速浏览原始日志"),
    ("告警中心", "聚焦异常\n深入分析", PURPLE, "风险等级筛选\n查看基线偏离证据\n查看 AI 研判结论"),
    ("用户基线", "对比画像\n量化偏离", ORANGE, "5W1H 六维画像\n每特征统计量\n小时活动分布图"),
    ("系统状态", "健康检查\n在线测试", GREEN, "5 服务状态卡\nKafka 消费延迟\n一键准确度测试"),
]
jw = Inches(2.82)
jgap = Inches(0.2)
jx = Inches(0.7)
for i, (title, purpose, clr, desc) in enumerate(journey):
    x = jx + i * (jw + jgap)
    _box(s, x, Inches(1.9), jw, Inches(2.1), fill=SUBTLE, border=BORDER, radius=0.06)
    _circle(s, x + Inches(1.1), Inches(2.05), Inches(0.35), fill=clr)
    _txt(s, x + Inches(1.1), Inches(2.1), Inches(0.35), Inches(0.28),
         str(i+1), size=13, color=BG, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, x + Inches(0.15), Inches(2.55), jw - Inches(0.3), Inches(0.25),
         title, size=13, color=clr, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, x + Inches(0.15), Inches(2.82), jw - Inches(0.3), Inches(0.22),
         purpose, size=10.5, color=WHITE, align=PP_ALIGN.CENTER)
    _line(s, x + Inches(0.3), Inches(3.1), jw - Inches(0.6), BORDER)
    _txt(s, x + Inches(0.15), Inches(3.2), jw - Inches(0.3), Inches(0.7),
         desc, size=9.5, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(journey) - 1:
        _arrow(s, x + jw, Inches(2.9), x + jw + jgap, Inches(2.9), color=DIM)

# Design principles
_txt(s, Inches(0.7), Inches(4.4), Inches(5.5), Inches(0.3),
     "前端架构设计思路", size=13, color=CYAN, bold=True)
principles = [
    ("API 层抽象", "apiFetch<T> 泛型封装，统一错误处理和类型推断。\n13 个 API 函数覆盖全部端点，前端不直接操作数据库。"),
    ("证据链可视化", "异常详情页不是简单展示字段，而是按因果链组织：\n告警摘要 → 基线偏离 → AI 研判 → 关联日志，逐层深挖。"),
    ("反馈闭环", "AI 研判结果包含 confidence 评分和 recommended_actions，\n分析师可据此决定是升级处理还是一键关闭。"),
]
y = Inches(4.8)
for title, desc in principles:
    _box(s, Inches(0.7), y, Inches(11.9), Inches(0.72), fill=SUBTLE, border=BORDER, radius=0.06)
    _txt(s, Inches(1.0), y + Inches(0.08), Inches(3), Inches(0.22),
         title, size=11.5, color=CYAN, bold=True)
    _txt(s, Inches(1.0), y + Inches(0.33), Inches(11.3), Inches(0.35),
         desc, size=10, color=GRAY)
    y += Inches(0.82)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — ACCURACY EVALUATION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 9, "检测效果评估：直接模式 vs 全链路对比", "ACCURACY EVALUATION")

# Big metrics
big_metrics = [
    ("53.5%", "F1 Score", "综合平衡\n精确率与召回率", CYAN),
    ("44.7%", "Precision", "告警中有多少\n是真正的异常", PURPLE),
    ("66.7%", "Recall", "真正的异常中\n有多少被检测到", ORANGE),
    ("300 条", "样本量", "3天×100条/天\n确定性种子生成", GREEN),
]
mw = Inches(2.75)
mgap = Inches(0.25)
mx = Inches(0.7)
for i, (val, label, desc, clr) in enumerate(big_metrics):
    x = mx + i * (mw + mgap)
    _box(s, x, Inches(1.5), mw, Inches(1.4), fill=SUBTLE, border=BORDER, radius=0.06)
    _txt(s, x + Inches(0.15), Inches(1.6), mw - Inches(0.3), Inches(0.5),
         val, size=28, color=clr, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, x + Inches(0.15), Inches(2.15), mw - Inches(0.3), Inches(0.25),
         label, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s, x + Inches(0.15), Inches(2.42), mw - Inches(0.3), Inches(0.38),
         desc, size=9, color=GRAY, align=PP_ALIGN.CENTER)

# Per dimension
_txt(s, Inches(0.7), Inches(3.2), Inches(5.5), Inches(0.3),
     "各维度表现差异与分析", size=14, color=CYAN, bold=True)

dim_data = [
    ("time", "时段", "80.0%", "P=100%", "R=66.7%", "非工作时段检测精准，小时直方图方法有效", GREEN),
    ("ip", "来源", "75.0%", "P=100%", "R=60.0%", "已知IP集合匹配准确，新IP判定逻辑可靠", GREEN),
    ("result", "结果", "82.4%", "P=87.5%", "R=77.8%", "登录成功率对比方法最稳定，历史数据充分", GREEN),
    ("access", "行为", "33.3%", "P=25%", "R=50.0%", "误报偏高—敏感资源加权过度，需调整阈值", ORANGE),
    ("geo", "地域", "0%", "—", "—", "缺少地理位置库数据，维度实际未生效", RED),
    ("volume", "流量", "0%", "—", "—", "窗口聚合已设计，实现延迟到后续版本", RED),
]

dim_w = Inches(5.75)
y = Inches(3.6)
for i in range(3):
    for j in range(2):
        idx = i * 2 + j
        if idx >= len(dim_data):
            continue
        dim_id, dim_cn, f1, prec, rec, note, clr = dim_data[idx]
        x = Inches(0.7) + j * (dim_w + Inches(0.25))
        yy = y + i * Inches(1.1)
        _box(s, x, yy, dim_w, Inches(0.95), fill=SUBTLE, border=BORDER, radius=0.06)
        _txt(s, x + Inches(0.15), yy + Inches(0.08), Inches(0.8), Inches(0.22),
             f"{dim_cn} ({dim_id})", size=11, color=clr, bold=True)
        _txt(s, x + Inches(1.2), yy + Inches(0.08), Inches(0.8), Inches(0.22),
             f"F1={f1}", size=11, color=WHITE, bold=True)
        _txt(s, x + Inches(2.2), yy + Inches(0.08), Inches(1.0), Inches(0.22),
             prec, size=10, color=GRAY)
        _txt(s, x + Inches(3.2), yy + Inches(0.08), Inches(1.0), Inches(0.22),
             rec, size=10, color=GRAY)
        _txt(s, x + Inches(0.15), yy + Inches(0.38), Inches(5.3), Inches(0.5),
             note, size=9.5, color=DIM)

# Evaluation method
_txt(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.35),
     "评估方法: 绕过 Kafka/ClickHouse，直接调用 RuleEngine + UebaScorer 对生成日志评分，与生成器内嵌的 ground truth 标签对比。确定性种子保证可复现。",
     size=10, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — KEY TECHNICAL DECISIONS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 10, "关键技术决策与取舍", "ARCHITECTURE DECISIONS")

decisions = [
    ("为什么 ClickHouse 而非 Elasticsearch?",
     "ES 擅长全文搜索，但安全日志分析的核心是聚合统计——\"这个用户过去7天的登录失败率均值是多少\""
     "这类查询在 ES 中需要复杂的聚合 DSL，在 ClickHouse 中是简单的 SELECT avg()。"
     "同时 ClickHouse 的 Kafka Engine 可以省掉一个消费者进程，数据直接从 Kafka 流入表。"
     "关联查询是弱项但通过多次查询 + API 层组合解决，满足当前场景。",
     "写入性能 3-5x, 存储压缩率 5-8x"),
    ("为什么 T+1 基线而非实时更新?",
     "实时更新基线意味着每次评分都要重算统计量，性能开销大且基线频繁波动。"
     "T+1 模式用昨天的数据描述今天的\"正常\"，稳定性好，计算成本低。"
     "对于大多数企业场景，用户行为模式在一天内不会有本质变化。",
     "简化计算, 稳定基线, 降低噪声"),
    ("为什么双引擎而非单一模型?",
     "纯 ML 模型在安全领域有可解释性问题——给一个 0.87 的风险分但说不出原因。"
     "纯规则覆盖面有限，无法检测未知攻击。两者互补: 规则保证已知威胁不丢，统计评分发现偏离。"
     "当前阶段以规则引擎为主 (命中即高危)，统计评分为辅 (作为佐证)，后续可逐步调整权重。",
     "可解释性 + 覆盖面平衡"),
    ("为什么 Flink + 独立消费者双路?",
     "Flink 负责流解析 (raw→parsed)，是 Flink 擅长的 ETL 工作；anomaly-detector 是 Python 进程，"
     "需要灵活调用基线缓存、数据库查询、AI API，用 PyFlink 实现反而不便。"
     "ClickHouse Kafka Engine 提供了第三条并行路径——不经过 Flink，直接从 Kafka 写入。",
     "职责分离, 各取所长"),
]

y = Inches(1.4)
for i, (question, answer, takeaway) in enumerate(decisions):
    _box(s, Inches(0.7), y, Inches(11.9), Inches(1.35), fill=SUBTLE, border=BORDER, radius=0.06)
    _txt(s, Inches(1.0), y + Inches(0.08), Inches(11.3), Inches(0.28),
         f"Q{i+1}: {question}", size=12.5, color=CYAN, bold=True)
    _txt(s, Inches(1.0), y + Inches(0.42), Inches(11.3), Inches(0.65),
         answer, size=10, color=GRAY)
    _txt(s, Inches(1.0), y + Inches(1.1), Inches(2), Inches(0.2),
         f"→ {takeaway}", size=9.5, color=GREEN)
    y += Inches(1.48)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — SUMMARY & ROADMAP
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_slide_title(s, 11, "当前成果与后续方向", "PROGRESS & ROADMAP")

# Left: Done
_txt(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.35),
     "已完成", size=15, color=GREEN, bold=True)
done_groups = [
    ("基础设施", "12 服务 Docker 编排 · 健康检查全覆盖\nFlink 守护进程自动监控 · 双路数据入库"),
    ("基线系统", "T+1 每日聚合 (23 指标) · 7 天滑动窗口\n5W1H 六维画像 · 多因子置信度评估\n新源持久化追踪 (ReplacingMergeTree)"),
    ("检测引擎", "6 条规则 (爆破/撞库/高频/越权/组合)\n六维自适应评分 · 双引擎合并策略"),
    ("AI 研判", "DeepSeek + DashScope 双模型容错\n5 层重试 + 指数退避 + pending 恢复\n结构化 JSON Prompt 工程"),
    ("前端 + API", "4 页面分析工作台 · 11 个 REST 端点\n准确度直接评估模式 (F1=53.5%)"),
]
y = Inches(2.0)
for group_title, items in done_groups:
    _txt(s, Inches(0.9), y, Inches(2), Inches(0.22),
         group_title, size=11, color=CYAN, bold=True)
    _txt(s, Inches(2.5), y, Inches(3.8), Inches(0.42),
         items, size=9.5, color=GRAY)
    y += Inches(0.48)

# Right: Next
_txt(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.35),
     "后续方向", size=15, color=ORANGE, bold=True)
next_items = [
    ("volume 窗口聚合", "实现延迟的窗口级别流量异常评分"),
    ("geo 地理位置", "集成 IP 地理位置库，激活地域维度"),
    ("access 调优", "降低敏感资源权重，减少误报"),
    ("同类群组 fallback", "低置信度用户使用部门/角色基线"),
    ("实时推送", "WebSocket/SSE 推送高危告警"),
    ("反馈闭环", "分析师标注 → 模型微调/规则更新"),
    ("多租户隔离", "tenant_id 维度完善"),
    ("生产压测", "scale profile 下的端到端性能基准"),
]
y = Inches(2.0)
for title, desc in next_items:
    _txt(s, Inches(7.2), y, Inches(2.2), Inches(0.22),
         f"→  {title}", size=11, color=ORANGE, bold=True)
    _txt(s, Inches(9.5), y, Inches(3.0), Inches(0.22),
         desc, size=9.5, color=GRAY)
    y += Inches(0.36)

# Bottom metrics bar
_bar = _box(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.2), fill=RGBColor(0x0C, 0x14, 0x1E))
_metrics_row = [
    ("12", "Docker\n服务"),
    ("6", "检测\n维度"),
    ("5", "重试\n层级"),
    ("11", "API\n端点"),
    ("4", "前端\n页面"),
    ("53.5%", "检测\nF1"),
]
for i, (num, label) in enumerate(_metrics_row):
    x = Inches(1.3) + i * Inches(2.05)
    _txt(s, x, Inches(6.0), Inches(1.8), Inches(0.42),
         num, size=22, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, x, Inches(6.42), Inches(1.8), Inches(0.4),
         label, size=9, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — END
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(s)
_box(s, 0, 0, Inches(0.12), H, fill=CYAN)
for i, clr in enumerate([CYAN, PURPLE, GREEN]):
    _circle(s, Inches(5.0) + i * Inches(0.4), Inches(2.2), Inches(0.2), fill=clr)
_txt(s, Inches(3), Inches(2.8), Inches(8), Inches(0.8),
     "谢谢", size=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
_txt(s, Inches(3), Inches(3.7), Inches(8), Inches(0.4),
     "UEBA 异常行为检测系统  |  阶段汇报", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── SAVE ───────────────────────────────────────────────────────
output = os.path.join(os.path.dirname(__file__), "..", "UEBA_项目阶段汇报.pptx")
prs.save(output)
print(f"PPT saved: {os.path.abspath(output)} ({os.path.getsize(output)//1024} KB)")
